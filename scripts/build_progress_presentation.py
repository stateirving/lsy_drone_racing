# ruff: noqa: ANN001, ANN201, E501
"""Build the SS26 progress presentation from the provided course template."""

from __future__ import annotations

import csv
import os
from pathlib import Path

os.environ.setdefault("MPLCONFIGDIR", "/tmp/matplotlib")

import matplotlib.pyplot as plt
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).parents[1]
TEMPLATE = ROOT / "Progress Presentation Template SS26.pptx"
OUTPUT = ROOT / "progress_presentation_ss26.pptx"
ASSET_DIR = ROOT / "presentation_assets"

TUM_BLUE = RGBColor(48, 112, 179)
TUM_DARK_BLUE = RGBColor(0, 82, 147)
LSY_GREEN = RGBColor(122, 184, 0)
LIGHT_BLUE = RGBColor(232, 241, 249)
LIGHT_GREEN = RGBColor(238, 247, 222)
LIGHT_GRAY = RGBColor(243, 245, 247)
MID_GRAY = RGBColor(112, 122, 130)
DARK = RGBColor(36, 45, 53)
ORANGE = RGBColor(227, 114, 34)
RED = RGBColor(196, 7, 27)
WHITE = RGBColor(255, 255, 255)


def delete_all_slides(prs: Presentation) -> None:
    """Remove the template example slides while keeping masters and layouts."""
    slide_ids = list(prs.slides._sldIdLst)  # noqa: SLF001
    for slide_id in slide_ids:
        prs.part.drop_rel(slide_id.rId)
        prs.slides._sldIdLst.remove(slide_id)  # noqa: SLF001


def rgb_hex(color: RGBColor) -> str:
    """Return a matplotlib-compatible hex color."""
    return f"#{color[0]:02x}{color[1]:02x}{color[2]:02x}"


def set_text(
    shape,
    text: str,
    *,
    size: float,
    color: RGBColor = DARK,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    margin: float = 0.06,
) -> None:
    """Set uniformly formatted text on a PowerPoint shape."""
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    p.line_spacing = 1.0
    for run in p.runs:
        run.font.name = "Arial"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color


def add_text(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    *,
    size: float = 14,
    color: RGBColor = DARK,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    margin: float = 0.02,
):
    """Add a formatted text box."""
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    set_text(
        shape,
        text,
        size=size,
        color=color,
        bold=bold,
        align=align,
        valign=valign,
        margin=margin,
    )
    return shape


def add_paragraphs(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    paragraphs: list[tuple[str, float, RGBColor, bool]],
    *,
    fill: RGBColor | None = None,
    line: RGBColor | None = None,
    radius: bool = True,
    margin: float = 0.13,
):
    """Add a text card with individually formatted paragraphs."""
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill or WHITE
    shape.line.color.rgb = line or (fill or WHITE)
    shape.line.width = Pt(1.1)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    for index, (text, size, color, bold) in enumerate(paragraphs):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = text
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(5 if index == 0 else 2)
        p.space_before = Pt(0)
        p.line_spacing = 1.0
        for run in p.runs:
            run.font.name = "Arial"
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
    return shape


def add_card(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    title: str,
    body: str,
    *,
    fill: RGBColor = WHITE,
    line: RGBColor = TUM_BLUE,
    title_color: RGBColor = TUM_DARK_BLUE,
    title_size: float = 13,
    body_size: float = 11.5,
    align: PP_ALIGN = PP_ALIGN.LEFT,
):
    """Add a standard title/body card."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.2)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.11)
    tf.margin_bottom = Inches(0.08)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    p.text = title
    p.alignment = align
    p.space_after = Pt(5)
    p.line_spacing = 1.0
    for run in p.runs:
        run.font.name = "Arial"
        run.font.size = Pt(title_size)
        run.font.bold = True
        run.font.color.rgb = title_color

    p = tf.add_paragraph()
    p.text = body
    p.alignment = align
    p.space_after = Pt(0)
    p.line_spacing = 1.0
    for run in p.runs:
        run.font.name = "Arial"
        run.font.size = Pt(body_size)
        run.font.color.rgb = DARK
    return shape


def add_pill(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    *,
    fill: RGBColor = TUM_BLUE,
    color: RGBColor = WHITE,
    size: float = 10,
):
    """Add a rounded label."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    set_text(
        shape,
        text,
        size=size,
        color=color,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        margin=0.01,
    )
    return shape


def add_chevron(slide, x: float, y: float, w: float, h: float, color: RGBColor = TUM_BLUE):
    """Add a small process arrow."""
    shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    return shape


def add_metric_card(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    value: str,
    label: str,
    *,
    accent: RGBColor = TUM_BLUE,
):
    """Add a large-number metric card."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = accent
    shape.line.width = Pt(1.5)
    compact = h < 0.9
    value_y = y + (0.04 if compact else 0.11)
    value_h = 0.28 if compact else 0.48
    value_size = 16 if compact else 25
    label_y = y + (0.35 if compact else 0.67)
    label_h = h - (0.38 if compact else 0.75)
    label_size = 8.2 if compact else 10.5
    add_text(
        slide,
        x + 0.1,
        value_y,
        w - 0.2,
        value_h,
        value,
        size=value_size,
        color=accent,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        x + 0.13,
        label_y,
        w - 0.26,
        label_h,
        label,
        size=label_size,
        color=DARK,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.TOP,
    )
    return shape


def add_picture_cover(slide, image_path: Path, x: float, y: float, w: float, h: float):
    """Add a picture cropped to cover a fixed rectangle."""
    with Image.open(image_path) as image:
        image_w, image_h = image.size
    image_ratio = image_w / image_h
    frame_ratio = w / h
    picture = slide.shapes.add_picture(
        str(image_path), Inches(x), Inches(y), width=Inches(w), height=Inches(h)
    )
    if image_ratio > frame_ratio:
        keep = frame_ratio / image_ratio
        crop = (1.0 - keep) / 2.0
        picture.crop_left = crop
        picture.crop_right = crop
    elif image_ratio < frame_ratio:
        keep = image_ratio / frame_ratio
        crop = (1.0 - keep) / 2.0
        picture.crop_top = crop
        picture.crop_bottom = crop
    return picture


def placeholder(slide, idx: int):
    """Return a placeholder by its layout index."""
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == idx:
            return shape
    raise KeyError(f"Placeholder {idx} not found on slide")


def set_header(slide, category: str, title: str, slide_no: int) -> None:
    """Populate the standard template header and footer."""
    set_text(placeholder(slide, 12), category.upper(), size=10, color=TUM_BLUE, bold=True)
    set_text(placeholder(slide, 0), title, size=25, color=DARK, bold=True)
    set_text(
        placeholder(slide, 11),
        f"Xiaochen Miao  |  Progress Presentation SS26  |  {slide_no}/5",
        size=8.5,
        color=MID_GRAY,
        align=PP_ALIGN.RIGHT,
        valign=MSO_ANCHOR.MIDDLE,
    )


def set_notes(slide, text: str) -> None:
    """Write presenter notes to the slide."""
    slide.notes_slide.notes_text_frame.text = text


def make_assets() -> tuple[Path, Path]:
    """Create the title still and checkpoint comparison chart."""
    ASSET_DIR.mkdir(exist_ok=True)

    title_image = ASSET_DIR / "race_title.png"
    with Image.open(ROOT / "docs/img/race.gif") as image:
        image.seek(0)
        image.convert("RGB").save(title_image)

    def read_series(path: Path) -> tuple[list[float], list[float]]:
        steps: list[float] = []
        success: list[float] = []
        with path.open(newline="") as file:
            for row in csv.DictReader(file):
                if row["step_m"] == "final":
                    continue
                steps.append(float(row["step_m"]))
                success.append(100.0 * float(row["success_rate"]))
        return steps, success

    smooth_steps, smooth_success = read_series(ROOT / "evaluation_level2_cmdtilt1p5_160M_summary.csv")
    fast_steps_a, fast_success_a = read_series(ROOT / "evaluation_level2_timepenalty001_summary.csv")
    fast_steps_b, fast_success_b = read_series(
        ROOT / "evaluation_level2_timepenalty001_from60_summary.csv"
    )
    fast_steps = fast_steps_a + fast_steps_b
    fast_success = fast_success_a + fast_success_b

    chart = ASSET_DIR / "checkpoint_success.png"
    fig, ax = plt.subplots(figsize=(8.2, 3.9), dpi=180)
    ax.plot(
        smooth_steps,
        smooth_success,
        marker="o",
        markersize=4,
        linewidth=2.0,
        color=rgb_hex(LSY_GREEN),
        label="Smoother command-tilt branch",
    )
    ax.plot(
        fast_steps,
        fast_success,
        marker="o",
        markersize=4,
        linewidth=2.0,
        color=rgb_hex(TUM_BLUE),
        label="Time-penalty branch",
    )
    ax.scatter([120], [87.5], s=75, color=rgb_hex(ORANGE), zorder=5)
    ax.annotate(
        "best: 87.5% @ 120M",
        xy=(120, 87.5),
        xytext=(76, 96),
        arrowprops={"arrowstyle": "->", "color": rgb_hex(ORANGE), "lw": 1.2},
        color=rgb_hex(DARK),
        fontsize=9,
        fontweight="bold",
    )
    ax.scatter([100], [70], s=70, facecolors="none", edgecolors=rgb_hex(DARK), linewidths=1.6)
    ax.annotate(
        "deployed smoother candidate",
        xy=(100, 70),
        xytext=(47, 59),
        arrowprops={"arrowstyle": "->", "color": rgb_hex(DARK), "lw": 1.0},
        color=rgb_hex(DARK),
        fontsize=8.5,
    )
    ax.set_xlim(5, 145)
    ax.set_ylim(-2, 102)
    ax.set_xlabel("Training steps [million]", fontsize=9)
    ax.set_ylabel("Success rate [%]", fontsize=9)
    ax.set_title("Fixed-seed checkpoint sweep: performance is not monotonic", fontsize=11, loc="left")
    ax.grid(True, alpha=0.25)
    ax.spines[["top", "right"]].set_visible(False)
    ax.tick_params(labelsize=8.5)
    ax.legend(frameon=False, fontsize=8.5, loc="upper left")
    fig.tight_layout()
    fig.savefig(chart, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return title_image, chart


def build_title_slide(prs: Presentation, title_image: Path) -> None:
    """Create the title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = placeholder(slide, 0)
    title.left, title.top, title.width, title.height = (
        Inches(0.7),
        Inches(1.45),
        Inches(6.1),
        Inches(1.55),
    )
    set_text(
        title,
        "Robust Level-2 Drone Racing\nwith Direct PPO",
        size=31,
        color=DARK,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )

    subtitle = placeholder(slide, 10)
    subtitle.left, subtitle.top, subtitle.width, subtitle.height = (
        Inches(0.72),
        Inches(3.25),
        Inches(5.8),
        Inches(0.38),
    )
    set_text(subtitle, "Progress Presentation  |  SS26", size=15, color=TUM_BLUE, bold=True)

    author = placeholder(slide, 11)
    author.left, author.top, author.width, author.height = (
        Inches(0.72),
        Inches(3.72),
        Inches(5.8),
        Inches(0.36),
    )
    set_text(author, "Xiaochen Miao", size=14, color=DARK)

    date = placeholder(slide, 12)
    date.left, date.top, date.width, date.height = (
        Inches(0.72),
        Inches(4.24),
        Inches(5.8),
        Inches(0.36),
    )
    set_text(date, "04 June 2026", size=12, color=MID_GRAY)

    panel = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.1), Inches(1.25), Inches(5.55), Inches(4.25)
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = LIGHT_BLUE
    panel.line.color.rgb = TUM_BLUE
    panel.line.width = Pt(1.2)
    add_picture_cover(slide, title_image, 7.23, 1.38, 5.29, 3.97)
    add_pill(slide, 9.16, 5.65, 3.05, 0.36, "CURRENT: LEVEL-2 POLICY", fill=LSY_GREEN, size=9)
    set_notes(
        slide,
        "Timing: 20 seconds.\n"
        "We are building a direct PPO controller for randomized Level-2 drone racing. "
        "The focus of this progress update is the method, the measured reliability gap, "
        "and the concrete plan to close it.",
    )


def build_method_overview(prs: Presentation) -> None:
    """Create methodology slide 1."""
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    set_header(slide, "Methodology / 1 of 3", "One policy closes the Level-2 loop", 1)

    box_x = [0.55, 3.02, 5.49, 7.96, 10.43]
    titles = ["OBSERVE", "ENCODE", "DECIDE", "COMMAND", "EXECUTE"]
    bodies = [
        "Drone state\nGates + obstacles",
        "Body-frame geometry\n2-step state history",
        "PPO actor\n2 x 64 tanh MLP",
        "Roll, pitch, yaw\nand thrust",
        "Crazyflow / Crazyflie\n50 Hz closed loop",
    ]
    fills = [LIGHT_BLUE, LIGHT_BLUE, LIGHT_GREEN, LIGHT_GREEN, LIGHT_BLUE]
    accents = [TUM_BLUE, TUM_BLUE, LSY_GREEN, LSY_GREEN, TUM_BLUE]
    for index, x in enumerate(box_x):
        add_card(
            slide,
            x,
            1.58,
            2.17,
            1.42,
            titles[index],
            bodies[index],
            fill=fills[index],
            line=accents[index],
            title_color=accents[index],
            title_size=11,
            body_size=11,
            align=PP_ALIGN.CENTER,
        )
        if index < 4:
            add_chevron(slide, x + 2.21, 2.08, 0.22, 0.38, TUM_BLUE)
    add_pill(slide, 5.18, 3.13, 3.0, 0.34, "NEW OBSERVATION EVERY 20 ms", fill=TUM_BLUE, size=9)

    add_paragraphs(
        slide,
        0.58,
        3.68,
        5.95,
        2.42,
        [
            ("Why direct PPO", 14, LSY_GREEN, True),
            ("One low-cost forward pass at runtime", 11.5, DARK, False),
            ("Reacts to newly observed geometry and dynamics", 11.5, DARK, False),
            ("Learns planning and control trade-offs jointly", 11.5, DARK, False),
            ("Selected for the randomized Level-2 task", 11.5, TUM_DARK_BLUE, True),
        ],
        fill=LIGHT_GREEN,
        line=LSY_GREEN,
    )
    add_paragraphs(
        slide,
        6.8,
        3.68,
        5.95,
        2.42,
        [
            ("Why not rely only on MPCC", 14, TUM_BLUE, True),
            ("Spline + OCP is interpretable and structured", 11.5, DARK, False),
            ("But randomized scenes require replanning and constraint tuning", 11.5, DARK, False),
            ("PPO limitation: reward-sensitive, with no hard collision guarantee", 11.5, RED, True),
        ],
        fill=LIGHT_BLUE,
        line=TUM_BLUE,
    )
    set_notes(
        slide,
        "Timing: 50 seconds.\n"
        "The controller uses a direct closed loop. Every 20 milliseconds it receives the current "
        "state and visible track geometry, converts it to body-frame features with short history, "
        "and produces attitude plus thrust. We selected PPO because online execution is just one "
        "network pass and it can react to randomization. Compared with the MPCC baseline, the cost "
        "is weaker guarantees and stronger dependence on reward design.",
    )


def build_observation_reward(prs: Presentation) -> None:
    """Create methodology slide 2."""
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    set_header(
        slide,
        "Methodology / 2 of 3",
        "Observation and reward design make racing learnable",
        2,
    )

    add_text(slide, 0.65, 1.45, 5.9, 0.35, "103-D body-frame observation", size=15, bold=True)
    groups = [
        ("State + progress", 17, TUM_BLUE),
        ("Gate geometry", 36, LSY_GREEN),
        ("Obstacles + visits", 20, ORANGE),
        ("Action + history", 30, RGBColor(116, 93, 160)),
    ]
    x = 0.65
    total_w = 5.95
    for label, dim, color in groups:
        width = total_w * dim / 103.0
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.93), Inches(width), Inches(0.63)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = WHITE
        set_text(
            shape,
            str(dim),
            size=12,
            color=WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
            margin=0.01,
        )
        x += width

    row_y = 2.76
    for index, (label, dim, color) in enumerate(groups):
        descriptions = [
            "height, body velocity, angular velocity, rotation, target progress",
            "previous, current and next gate corners in the drone frame",
            "relative obstacle positions plus visited gate/obstacle flags",
            "previous normalized action plus two 13-D state-history samples",
        ]
        add_card(
            slide,
            0.65,
            row_y + index * 0.69,
            5.95,
            0.58,
            f"{label}  ({dim}D)",
            descriptions[index],
            fill=WHITE,
            line=color,
            title_color=color,
            title_size=10.5,
            body_size=9.4,
        )

    add_card(
        slide,
        0.65,
        5.69,
        5.95,
        0.62,
        "4-D normalized action",
        "roll | pitch | yaw | thrust  ->  scaled to the attitude controller limits",
        fill=LIGHT_GREEN,
        line=LSY_GREEN,
        title_color=LSY_GREEN,
        title_size=10.5,
        body_size=9.7,
    )

    add_text(slide, 7.02, 1.45, 5.65, 0.35, "Shaped reward follows a gate-passing sequence", size=15, bold=True)
    stage_x = [7.02, 8.47, 9.92, 11.37]
    stage_titles = ["APPROACH", "CENTER", "PASS", "EXIT"]
    stage_bodies = ["axis progress", "gate stage", "pass bonus", "back-gate bonus"]
    for index, x in enumerate(stage_x):
        add_card(
            slide,
            x,
            1.94,
            1.22,
            1.05,
            stage_titles[index],
            stage_bodies[index],
            fill=LIGHT_BLUE if index < 2 else LIGHT_GREEN,
            line=TUM_BLUE if index < 2 else LSY_GREEN,
            title_color=TUM_BLUE if index < 2 else LSY_GREEN,
            title_size=9.2,
            body_size=9,
            align=PP_ALIGN.CENTER,
        )
        if index < 3:
            add_chevron(slide, x + 1.26, 2.27, 0.14, 0.3, TUM_BLUE)

    add_paragraphs(
        slide,
        7.02,
        3.28,
        2.72,
        2.18,
        [
            ("Positive signals", 13, LSY_GREEN, True),
            ("Gate-axis / stage progress", 10.8, DARK, False),
            ("Gate pass + finish bonus", 10.8, DARK, False),
            ("Fast completion", 10.8, DARK, False),
        ],
        fill=LIGHT_GREEN,
        line=LSY_GREEN,
    )
    add_paragraphs(
        slide,
        9.95,
        3.28,
        2.72,
        2.18,
        [
            ("Safety + smoothness costs", 13, RED, True),
            ("Crash / obstacle proximity", 10.8, DARK, False),
            ("Tilt and command tilt", 10.8, DARK, False),
            ("Action changes + time", 10.8, DARK, False),
        ],
        fill=RGBColor(252, 238, 240),
        line=RED,
    )
    add_pill(
        slide,
        7.02,
        5.76,
        5.65,
        0.42,
        "BODY-FRAME GEOMETRY REDUCES GLOBAL-COORDINATE DEPENDENCE",
        fill=TUM_DARK_BLUE,
        size=8.8,
    )
    set_notes(
        slide,
        "Timing: 55 seconds.\n"
        "The policy receives 103 features. The important choice is to express gate and obstacle "
        "geometry in the drone body frame and include previous, current, and next gates plus short "
        "state history. Sparse race completion alone was not enough, so the reward decomposes gate "
        "passing into approach, centering, crossing, and exiting. Safety, tilt, smoothness, and time "
        "penalties create the speed-versus-robustness trade-off.",
    )


def build_training_diagnosis(prs: Presentation, chart: Path) -> None:
    """Create methodology slide 3."""
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    set_header(slide, "Methodology / 3 of 3", "Training is paired with reproducible diagnosis", 3)

    chart_panel = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(1.5), Inches(7.65), Inches(4.37)
    )
    chart_panel.fill.solid()
    chart_panel.fill.fore_color.rgb = WHITE
    chart_panel.line.color.rgb = RGBColor(210, 218, 224)
    add_picture_cover(slide, chart, 0.68, 1.63, 7.39, 4.08)

    pipeline = [
        ("1  Parallel simulation", "1,024 Crazyflow environments on GPU", TUM_BLUE),
        ("2  PPO update", "32-step rollouts, GAE, actor-critic optimization", LSY_GREEN),
        ("3  Checkpoint sweep", "Save every 10M steps; evaluate deterministic policy", ORANGE),
        ("4  Failure diagnosis", "Fixed seeds + crash-object and hotspot analysis", RED),
    ]
    for index, (title, body, color) in enumerate(pipeline):
        add_card(
            slide,
            8.48,
            1.5 + index * 1.12,
            4.28,
            0.96,
            title,
            body,
            fill=WHITE,
            line=color,
            title_color=color,
            title_size=11.3,
            body_size=10,
        )

    add_metric_card(slide, 8.48, 5.98, 2.05, 0.66, "150M", "training steps", accent=TUM_BLUE)
    add_metric_card(slide, 10.71, 5.98, 2.05, 0.66, "2h 26m", "GPU training time", accent=LSY_GREEN)
    set_notes(
        slide,
        "Timing: 55 seconds.\n"
        "Training uses 1,024 parallel environments, 32-step rollouts, and regular checkpoints. "
        "A 150-million-step run took about two hours and 26 minutes on GPU. We do not select a "
        "policy from training reward alone. Each checkpoint is evaluated deterministically on "
        "fixed seeds, then crashes are classified by target gate, object, and location. The curve "
        "shows why this matters: success is not monotonic with more training.",
    )


def build_progress(prs: Presentation) -> None:
    """Create the progress-to-date slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    set_header(slide, "Progress to date", "Fast policy achieved; reliability is the remaining gap", 4)

    metric_data = [
        ("87.5%", "best success rate\n40 fixed-seed episodes", TUM_BLUE),
        ("5.80 s", "mean successful time\nbest 120M checkpoint", LSY_GREEN),
        ("0%", "position safety-limit violations\nacross the 40-seed sweep", TUM_DARK_BLUE),
        ("72.5%", "smoother deployed candidate\n200 single-env episodes", ORANGE),
    ]
    for index, (value, label, color) in enumerate(metric_data):
        add_metric_card(slide, 0.55 + index * 3.12, 1.5, 2.91, 1.27, value, label, accent=color)

    add_text(slide, 0.58, 2.93, 8.35, 0.27, "Crash hotspots from 1,024 randomized vector worlds", size=11.5, bold=True)
    hotspot = ROOT / "evaluation_level2_100M_crash_hotspots_hotspots.png"
    add_picture_cover(slide, hotspot, 0.55, 3.25, 8.48, 2.67)

    add_paragraphs(
        slide,
        9.25,
        3.25,
        3.52,
        2.67,
        [
            ("What did not work as expected", 13, RED, True),
            ("Success regresses after some checkpoints; more training is not automatically better.", 10.5, DARK, False),
            ("318 / 1,024 crashes: Gate 3 top = 132; Obstacle 2 = 96.", 10.5, DARK, False),
            ("Faster policies produce more aggressive tilt commands.", 10.5, DARK, False),
            ("Response: targeted hotspot diagnosis + checkpoint selection.", 10.5, TUM_DARK_BLUE, True),
        ],
        fill=RGBColor(252, 245, 245),
        line=RED,
    )
    add_pill(
        slide,
        9.52,
        6.08,
        2.96,
        0.37,
        "NEXT: FIX TWO DOMINANT HOTSPOTS",
        fill=RED,
        size=8.5,
    )
    set_notes(
        slide,
        "Timing: 70 seconds.\n"
        "The best checkpoint currently reaches 87.5 percent success over 40 fixed seeds with a "
        "5.80-second mean successful race and no position-limit violations. For deployment we are "
        "currently using a smoother 100-million-step policy, which reaches 72.5 percent over 200 "
        "single-environment episodes. The remaining gap is concentrated: in a 1,024-world test, "
        "Gate 3 top and Obstacle 2 account for 228 of 318 crashes. More training is not monotonic, "
        "so the next iteration targets those failure modes instead of simply increasing steps.",
    )


def build_milestones(prs: Presentation) -> None:
    """Create the milestones and SMART goals slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    set_header(slide, "Milestones and SMART goals", "Plan for the remainder of the semester", 5)

    add_pill(
        slide,
        0.55,
        1.45,
        12.22,
        0.55,
        "OBJECTIVE BY 02 JULY: RELIABLE LEVEL-2 POLICY, VALIDATED AND DEPLOYMENT-READY",
        fill=TUM_DARK_BLUE,
        size=10.5,
    )

    cards = [
        (
            "11 JUN",
            "Hotspot fix",
            "Change reward / observation around Gate 3 top and Obstacle 2.\n\n"
            "Acceptance: 1,024-episode test; reduce these two crash counts by >=30%.",
            RED,
        ),
        (
            "18 JUN",
            "Policy freeze",
            "Compare candidate checkpoints on the same randomized worlds.\n\n"
            "Acceptance: >=85% success, <=6.0 s mean successful time, 0 position-limit violations.",
            TUM_BLUE,
        ),
        (
            "25 JUN",
            "Sim-to-real validation",
            "Run the frozen policy through the deployment pipeline and lab track.\n\n"
            "Acceptance: at least 18 / 20 completed runs and 0 safety-stop events.",
            LSY_GREEN,
        ),
        (
            "02 JUL",
            "Final package",
            "Freeze checkpoint, config, evaluation outputs, and final report material.\n\n"
            "Acceptance: one-command reproducible evaluation and deployment-ready controller.",
            ORANGE,
        ),
    ]
    for index, (date, title, body, color) in enumerate(cards):
        x = 0.5 + index * 3.18
        add_pill(slide, x + 0.17, 2.26, 2.38, 0.38, date, fill=color, size=9.5)
        add_card(
            slide,
            x,
            2.48,
            2.73,
            3.55,
            title,
            body,
            fill=WHITE,
            line=color,
            title_color=color,
            title_size=13,
            body_size=10.2,
        )
        add_pill(
            slide,
            x + 0.3,
            5.55,
            2.13,
            0.32,
            "OWNER: XIAOCHEN MIAO",
            fill=LIGHT_GRAY,
            color=DARK,
            size=7.8,
        )
        if index < 3:
            add_chevron(slide, x + 2.82, 4.0, 0.25, 0.42, TUM_BLUE)

    add_text(
        slide,
        0.62,
        6.23,
        12.0,
        0.32,
        "Risk buffer: retain the 120M time-penalty checkpoint as the fallback if hotspot tuning or hardware validation regresses.",
        size=9.5,
        color=RED,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    set_notes(
        slide,
        "Timing: 50 seconds.\n"
        "The remaining work is split into four measurable milestones. First, reduce the two dominant "
        "hotspot crash counts by at least 30 percent in a 1,024-episode test. Second, freeze a policy "
        "that reaches at least 85 percent success, stays below six seconds, and has no position-limit "
        "violations. Third, complete at least 18 of 20 lab runs without a safety stop. Finally, freeze "
        "the reproducible evaluation and deployment package. Xiaochen Miao owns each task in the "
        "current repository workflow.",
    )


def validate(prs: Presentation) -> None:
    """Fail fast on common deck generation errors."""
    if len(prs.slides) != 6:
        raise RuntimeError(f"Expected 6 slides, got {len(prs.slides)}")
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    for slide_index, slide in enumerate(prs.slides, start=1):
        if not any(getattr(shape, "text", "").strip() for shape in slide.shapes):
            raise RuntimeError(f"Slide {slide_index} contains no text")
        for shape in slide.shapes:
            if shape.left < 0 or shape.top < 0:
                raise RuntimeError(f"Slide {slide_index}: negative shape position for {shape.name}")
            if shape.left + shape.width > slide_w + Inches(0.02):
                raise RuntimeError(f"Slide {slide_index}: shape exceeds slide width: {shape.name}")
            if shape.top + shape.height > slide_h + Inches(0.02):
                raise RuntimeError(f"Slide {slide_index}: shape exceeds slide height: {shape.name}")


def main() -> None:
    """Build and save the presentation."""
    if not TEMPLATE.exists():
        raise FileNotFoundError(f"Missing template: {TEMPLATE}")
    title_image, chart = make_assets()
    prs = Presentation(TEMPLATE)
    delete_all_slides(prs)
    build_title_slide(prs, title_image)
    build_method_overview(prs)
    build_observation_reward(prs)
    build_training_diagnosis(prs, chart)
    build_progress(prs)
    build_milestones(prs)
    validate(prs)
    prs.save(OUTPUT)
    print(f"saved {OUTPUT}")


if __name__ == "__main__":
    main()
